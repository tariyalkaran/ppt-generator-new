# slide_extractor.py
import os
import tempfile
import uuid
from copy import deepcopy
from pptx import Presentation
from pptx.util import Inches, Pt
from PIL import Image, ImageDraw, ImageFont
from utils import get_env, logger
from azure.storage.blob import BlobServiceClient

AZURE_CONN = get_env("AZURE_BLOB_CONN", required=True)
BLOB_CONTAINER = get_env("AZURE_BLOB_CONTAINER", "ppt-dataset")

def download_blob_to_local(blob_name: str, dest_path: str):
    """
    Download blob from the source container to a local file path.
    """
    try:
        blob_service = BlobServiceClient.from_connection_string(AZURE_CONN)
        container_client = blob_service.get_container_client(BLOB_CONTAINER)
        with open(dest_path, "wb") as fp:
            stream = container_client.download_blob(blob_name)
            stream.readinto(fp)
        return dest_path
    except Exception as e:
        logger.exception(f"Failed to download blob {blob_name}: {e}")
        raise

def extract_slides_info_from_ppt(local_ppt_path: str):
    """
    Return list of slide metadata dicts:
    {
      "slide_index": int,
      "title": str,
      "text": str,
      "preview_image": "/tmp/...",
      "ppt_path": local_ppt_path,
      "slide_id": "<pptbasename>_Slide_XX"
    }
    """
    prs = Presentation(local_ppt_path)
    slides_info = []
    base = os.path.splitext(os.path.basename(local_ppt_path))[0]

    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text and shape.text.strip():
                texts.append(shape.text.strip())
        title = texts[0] if texts else f"Slide {i+1}"
        combined_text = "\n".join(texts)
        slide_id = f"{base}_Slide_{i:02d}"

        preview_image = _make_text_preview_image(title, combined_text)

        slides_info.append({
            "slide_index": i,
            "title": title,
            "text": combined_text,
            "preview_image": preview_image,
            "ppt_path": local_ppt_path,
            "slide_id": slide_id
        })

    return slides_info

def _make_text_preview_image(title: str, body_text: str, width=800, height=450):
    """
    Create a simple preview PNG showing the title and first few bullet lines.
    This is for UI selection only.
    """
    try:
        img = Image.new("RGB", (width, height), color=(245, 246, 250))
        draw = ImageDraw.Draw(img)

        # attempt to use a default font; fallback to PIL default
        try:
            font_title = ImageFont.truetype("DejaVuSans-Bold.ttf", 24)
            font_body = ImageFont.truetype("DejaVuSans.ttf", 16)
        except Exception:
            font_title = ImageFont.load_default()
            font_body = ImageFont.load_default()

        padding = 20
        y = padding

        # Title
        draw.text((padding, y), title, font=font_title, fill=(0, 85, 170))
        y += 40

        # Body: show first 6 lines
        lines = body_text.splitlines()[:6]
        for ln in lines:
            # simple bullet formatting
            draw.text((padding + 10, y), u"\u2022 " + ln[:120], font=font_body, fill=(40, 40, 40))
            y += 22

        tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".png")
        img.save(tmp.name, format="PNG")
        tmp.close()
        return tmp.name
    except Exception as e:
        logger.exception(f"Failed to create preview image: {e}")
        return None