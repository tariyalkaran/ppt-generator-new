# =============================================
# generate_ppt_llm.py
# =============================================
import os
import uuid
from pptx import Presentation
from pptx.util import Pt
from utils import text_client, get_env, logger


# ------------------------------------------------------------
# TITLE SLIDE
# ------------------------------------------------------------
def add_title_slide(prs, title_text, subtitle_text=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank slide

    title_box = slide.shapes.add_textbox(
        left=prs.slide_width * 0.15,
        top=prs.slide_height * 0.35,
        width=prs.slide_width * 0.7,
        height=Pt(80),
    )

    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(40)
    p.font.bold = True
    p.alignment = 1  # CENTER

    if subtitle_text:
        sub_box = slide.shapes.add_textbox(
            left=prs.slide_width * 0.15,
            top=prs.slide_height * 0.48,
            width=prs.slide_width * 0.7,
            height=Pt(40),
        )
        tf2 = sub_box.text_frame
        tf2.clear()
        sp = tf2.paragraphs[0]
        sp.text = subtitle_text
        sp.font.size = Pt(18)
        sp.alignment = 1  # CENTER


# ------------------------------------------------------------
# LLM BULLET GENERATOR (UNCHANGED)
# ------------------------------------------------------------
def llm_synthesize_slide(user_answers, global_prompt):
    qa_text = "\n".join(
        f"Q: {q}\nA: {a}"
        for q, a in user_answers.items()
        if a.strip()
    )

    prompt = f"""
You are a senior consultant creating a professional PowerPoint slide.

GLOBAL CONTEXT:
{global_prompt}

USER INPUT:
{qa_text}

TASK:
- derive a slide title
- derive 5–6 bullets
- Paraphrase and enrich & rewrite clearly
- professional business tone

FORMAT:
Title: <title>
- bullet
- bullet
"""

    resp = text_client.chat.completions.create(
        model=get_env("CHAT_MODEL", required=True),
        messages=[{"role": "user", "content": prompt}],
        max_tokens=500,
        temperature=0.7,
    )

    raw = resp.choices[0].message.content.strip()
    lines = [x.strip() for x in raw.split("\n") if x.strip()]

    title = "Slide"
    bullets = []

    for ln in lines:
        if ln.lower().startswith("title"):
            title = ln.split(":", 1)[1].strip()
        elif ln.startswith("-"):
            bullets.append(ln[1:].strip())

    return title, bullets


# ============================================================
# MAIN PPT GENERATOR (PREVIEW-SAFE, FINAL)
# ============================================================
def generate_presentation(payload):
    slides = payload.get("preview_slides") or payload.get("slides", [])
    answers_map = payload.get("answers_map", {})

    if not slides:
        raise ValueError("No slides provided")

    prs = Presentation()
    global_prompt = "professional business presentation"

    # ---------------------------------------------------
    # 🔑 DETECT PREVIEW MODE
    # ---------------------------------------------------
    preview_mode = (
        isinstance(slides, list)
        and slides
        and isinstance(slides[0], dict)
        and "title" in slides[0]
        and "bullets" in slides[0]
        and "slide_index" not in slides[0]
    )

    # ===================================================
    # PREVIEW MODE → DIRECT PPT GENERATION
    # ===================================================
    if preview_mode:
        for slide in slides:
            title = (slide.get("title") or "").strip()
            bullets = slide.get("bullets") or []

            # ✅ Normalize bullets (VERY IMPORTANT)
            bullets = [
                b.strip()
                for b in bullets
                if isinstance(b, str) and b.strip()
            ]

            # -----------------------------
            # Title-only slide
            # -----------------------------
            if not bullets:
                add_title_slide(prs, title)
                continue

            # -----------------------------
            # Title + bullets slide
            # -----------------------------
            ppt_slide = prs.slides.add_slide(prs.slide_layouts[1])
            ppt_slide.shapes.title.text = title

            body = ppt_slide.placeholders[1].text_frame
            body.clear()

            for b in bullets:
                para = body.add_paragraph()
                para.text = b
                para.level = 0
                para.font.size = Pt(18)

    # ===================================================
    # ORIGINAL Q&A → LLM → PPT FLOW (UNCHANGED)
    # ===================================================
    else:
        for slide in slides:
            slide_idx = str(slide["slide_index"])
            slide_title = slide["slide_title"]
            user_answers = answers_map.get(slide_idx, {})

            # -----------------------------
            # CASE 1 — Presentation title slide
            # -----------------------------
            if "What should be the title of this presentation?" in user_answers:
                title_text = user_answers[
                    "What should be the title of this presentation?"
                ].strip()

                subtitle = None
                for v in user_answers.values():
                    if "202" in v:
                        subtitle = v.strip()
                        break

                add_title_slide(prs, title_text, subtitle)
                continue

            ppt_slide = prs.slides.add_slide(prs.slide_layouts[1])

            # -----------------------------
            # CASE 2 — User skipped answers
            # -----------------------------
            valid_answers_exist = any(v.strip() for v in user_answers.values())
            if not valid_answers_exist:
                ppt_slide.shapes.title.text = slide_title
                continue

            # -----------------------------
            # CASE 3 — Normal LLM generation
            # -----------------------------
            try:
                _, bullets = llm_synthesize_slide(
                    user_answers,
                    global_prompt
                )
            except Exception:
                logger.exception("LLM failed")
                bullets = []

            ppt_slide.shapes.title.text = slide_title
            body = ppt_slide.placeholders[1].text_frame
            body.clear()

            for b in bullets:
                para = body.add_paragraph()
                para.text = b
                para.level = 0
                para.font.size = Pt(18)

    # ---------------------------------------------------
    # OPTIONAL ENDING SLIDE
    # ---------------------------------------------------
    thanks = prs.slides.add_slide(prs.slide_layouts[1])
    thanks.shapes.title.text = "Thank You"

    # ---------------------------------------------------
    # SAVE FILE
    # ---------------------------------------------------
    os.makedirs("generated", exist_ok=True)
    out_path = f"generated/ppt_{uuid.uuid4().hex[:6]}.pptx"
    prs.save(out_path)
    return out_path