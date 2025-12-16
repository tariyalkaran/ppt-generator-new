# slide_renderer.py
import os
import uuid
import pythoncom
import win32com.client
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def export_slide_to_png(ppt_path, slide_index):
    pythoncom.CoInitialize()
    powerpoint = win32com.client.Dispatch("PowerPoint.Application")
    powerpoint.Visible = True

    pres = powerpoint.Presentations.Open(ppt_path, WithWindow=True)
    slide = pres.Slides[slide_index]

    out_path = os.path.join(
        os.path.dirname(ppt_path),
        f"slide_{slide_index}_{uuid.uuid4().hex[:6]}.png"
    )

    slide.Export(out_path, "PNG", 1920, 1080)

    pres.Close()
    powerpoint.Quit()
    pythoncom.CoUninitialize()

    return out_path


def _is_editable_text_shape(shape):
    if not shape.has_text_frame:
        return False

    text = shape.text.strip()
    if not text:
        return False

    if len(text) < 3:
        return False

    return True


def extract_slide_structure(ppt_path, slide_index):
    prs = Presentation(ppt_path)
    slide = prs.slides[slide_index]

    editable_shapes = []
    idx = 0

    for shape in slide.shapes:

        # -----------------------------
        # TEXT SHAPES
        # -----------------------------
        if _is_editable_text_shape(shape):

            # ✅ SAFE placeholder detection
            is_placeholder = False
            try:
                _ = shape.placeholder_format
                is_placeholder = True
            except Exception:
                is_placeholder = False

            shape_entry = {
                "shape_id": f"shape_{idx}",
                "text": shape.text.strip(),
                "placeholder": is_placeholder,
                "type": "title" if is_placeholder else "body"
            }

            editable_shapes.append(shape_entry)
            idx += 1

        # -----------------------------
        # GROUP SHAPES
        # -----------------------------
        elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for shp in shape.shapes:
                if _is_editable_text_shape(shp):
                    editable_shapes.append({
                        "shape_id": f"shape_{idx}",
                        "text": shp.text.strip(),
                        "placeholder": False,
                        "type": "body"
                    })
                    idx += 1

    png_path = export_slide_to_png(ppt_path, slide_index)

    return {
        "slide_index": slide_index,
        "ppt_path": ppt_path,
        "png_path": png_path,
        "editable_shapes": editable_shapes
    }