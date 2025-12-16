import os 
import uuid
import tempfile
from pptx import Presentation
from azure.storage.blob import BlobServiceClient
from openai import AzureOpenAI
from chromadb import PersistentClient
from utils import get_env, logger, now_ts, get_embedding_dim

# === CONFIG ===
BLOB_CONN = get_env("AZURE_BLOB_CONN", required=True)
BLOB_CONTAINER = get_env("AZURE_BLOB_CONTAINER", "ppt-dataset")
EMBEDDING_MODEL = get_env("EMBEDDING_MODEL", "text-embedding-3-small")
CHROMA_PERSIST_DIR = get_env("CHROMA_PERSIST_DIR", "./chroma_db")

# === AZURE OPENAI CLIENT ===
text_client = AzureOpenAI(
    azure_endpoint=get_env("OPENAI_API_BASE", required=True),
    api_key=get_env("OPENAI_API_KEY", required=True),
    api_version=get_env("OPENAI_API_VERSION", "2024-05-01-preview")
)

# === AZURE BLOB CLIENT ===
blob_client = BlobServiceClient.from_connection_string(BLOB_CONN)
container_client = blob_client.get_container_client(BLOB_CONTAINER)

# === CHROMA CLIENT (new syntax) ===
chroma_client = PersistentClient(path=CHROMA_PERSIST_DIR)
try:
    collection = chroma_client.get_collection("ppt_slides")
except Exception:
    collection = chroma_client.create_collection("ppt_slides")

EMBEDDING_DIM = get_embedding_dim(EMBEDDING_MODEL)

# === FUNCTIONS ===
def extract_slides(local_path):
    """Extract text content from all slides in a PPT."""
    prs = Presentation(local_path)
    slides = []
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text and shape.text.strip():
                texts.append(shape.text.strip())
        slides.append({"index": i, "text": "\n".join(texts)})
    return slides


def simple_tagger(text):
    """Create simple topic tags based on content keywords."""
    text_l = text.lower()
    tags = set()
    if any(k in text_l for k in ["design", "architecture", "ui", "ux"]):
        tags.add("Design")
    if any(k in text_l for k in ["test", "qa", "verification"]):
        tags.add("Test")
    if any(k in text_l for k in ["migration", "migrate"]):
        tags.add("Migration")
    for domain in ["claims", "membership", "provider", "finance", "medicaid", "commercial"]:
        if domain in text_l:
            tags.add(domain.capitalize())
    return list(tags) or ["General"]


def ppt_already_indexed(ppt_name):
    """Check if PPT is already embedded in Chroma."""
    try:
        res = collection.query(where={"ppt_name": ppt_name}, n_results=1)
        ids = res.get("ids", [[]])[0]
        return len(ids) > 0
    except Exception:
        return False


def azure_embed_func(texts):
    """Generate embeddings from Azure OpenAI."""
    try:
        resp = text_client.embeddings.create(model=EMBEDDING_MODEL, input=texts)
        return [d.embedding for d in resp.data]
    except Exception as e:
        logger.exception(f"Embedding failed: {e}")
        return []


def process_blob(blob_name):
    """Download PPT, extract slides, generate embeddings, and insert into Chroma."""
    logger.info(f"Processing blob: {blob_name}")
    tmp_path = os.path.join(tempfile.gettempdir(), blob_name.replace("/", "_"))
    with open(tmp_path, "wb") as fp:
        stream = container_client.download_blob(blob_name)
        stream.readinto(fp)

    slides = extract_slides(tmp_path)
    if not slides:
        logger.warning(f"No slides found in {blob_name}")
        return

    if ppt_already_indexed(blob_name):
        logger.info(f"Skipping '{blob_name}' — already indexed in Chroma.")
        return

    docs, metadatas, ids, texts = [], [], [], []
    for s in slides:
        slide_id = f"{os.path.splitext(os.path.basename(blob_name))[0]}_Slide_{s['index']:02d}"
        text = s.get("text", "") or ""
        metadata = {
            "ppt_name": blob_name,
            "slide_index": str(s["index"]),
            "slide_id": slide_id,
            "title": text.split("\n", 1)[0] if text else "",
            "tags": ", ".join(simple_tagger(text)), # ✅ FIXED: Convert list → string
            "indexed_on": str(now_ts()) # ✅ FIXED: Ensure string
        }
        ids.append(str(uuid.uuid4()))
        docs.append(text)
        metadatas.append(metadata)
        texts.append(text)

    embeddings = azure_embed_func(texts)
    if not embeddings or len(embeddings) != len(docs):
        logger.error("Embedding count mismatch or failed; aborting indexing for this file.")
        return

    # ✅ Insert into Chroma
    try:
        collection.add(documents=docs, embeddings=embeddings, metadatas=metadatas, ids=ids)
        logger.info(f"Indexed {len(docs)} slides from {blob_name} into Chroma.")
    except Exception as e:
        logger.exception(f"Failed to insert slides from {blob_name} into Chroma: {e}")


def delete_ppt_from_chroma(ppt_name: str) -> None:
    """
    Delete all Chroma slide indexes that belong to the given PPT.
    Matches your metadata:
    metadata = { "ppt_name": blob_name, ... }
    """

    logger.info(f"Deleting Chroma indexes for PPT: {ppt_name}")

    try:
        # ✅ DIRECT DELETE — no pre-query (avoids Chroma API bug)
        collection.delete(where={"ppt_name": ppt_name})

        # ✅ IMPORTANT: Persist the deletion to disk
        # chroma_client.persist()

        logger.info(f"✅ Successfully deleted Chroma indexes for PPT: {ppt_name}")

    except Exception as e:
        logger.exception(f"❌ Failed to delete Chroma indexes for PPT: {ppt_name}")
        raise e


def main():
    """Main ingestion process."""
    logger.info("Starting ingestion into Chroma from Azure Blob (ppt-dataset)...")
    for b in container_client.list_blobs():
        if b.name.endswith(".pptx") or b.name.endswith(".ppt"):
            try:
                process_blob(b.name)
            except Exception as e:
                logger.exception(f"Failed to process {b.name}: {e}")

    logger.info("Ingestion complete.")


if __name__ == "__main__":
    main()
