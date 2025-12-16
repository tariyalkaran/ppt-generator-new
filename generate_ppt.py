# generate_ppt.py
import os
import uuid
from pptx import Presentation
from pptx.util import Pt
from utils import logger

def add_title_slide(prs, title):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = "Prepared by AI PPT Generator"

def add_agenda(prs, titles):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Agenda"
    body = slide.placeholders[1].text_frame
    body.clear()
    for t in titles:
        body.add_paragraph().text = t

def add_content_slide(prs, title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for b in bullets:
        tf.add_paragraph().text = b

def add_thankyou(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Thank You"

def generate_presentation(payload):
    """
    payload = {
        "selected_slides": [slide_structs...],
        "answers_map": { slide_index: {shape_id: text} }
    }
    """

    selected_slides = payload.get("selected_slides", [])
    answers_map = payload.get("answers_map", {})

    if not selected_slides:
        raise ValueError("No slides selected")

    prs = Presentation()

    # ---------- TITLE SLIDE ----------
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = "Generated Presentation"
    subtitle = title_slide.placeholders[1]
    subtitle.text = "Auto-generated using AI"

    # ---------- CONTENT SLIDES ----------
    for idx, slide in enumerate(selected_slides, start=1):

        slide_title = slide.get("title", f"Slide {idx}")
        slide_id = slide.get("slide_id")
        slide_answers = answers_map.get(str(slide["slide_index"]), {})

        content_slide = prs.slides.add_slide(prs.slide_layouts[1])
        content_slide.shapes.title.text = slide_title

        body = content_slide.shapes.placeholders[1].text_frame
        body.clear()

        if slide_answers:
            for _, text in slide_answers.items():
                p = body.add_paragraph()
                p.text = text.strip()
                p.font.size = Pt(18)
                p.level = 0
        else:
            p = body.add_paragraph()
            p.text = "Content to be updated"
            p.level = 0

    # ---------- THANK YOU SLIDE ----------
    thank_slide = prs.slides.add_slide(prs.slide_layouts[1])
    thank_slide.shapes.title.text = "Thank You"
    thank_slide.shapes.placeholders[1].text = "Questions?"

    # ---------- SAVE ----------
    os.makedirs("generated", exist_ok=True)
    out_path = f"generated/ppt_{uuid.uuid4().hex[:6]}.pptx"
    prs.save(out_path)

    return out_path