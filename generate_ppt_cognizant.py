# ============================================================
# generate_ppt_cognizant.py
# Cognizant Template | Preview-driven | TEXT-ONLY
# ============================================================

import os
import uuid
from datetime import datetime
from copy import deepcopy

from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor

from utils import logger


BASE_DIR = os.path.dirname(os.path.abspath(__file__))
TEMPLATES_DIR = os.path.join(BASE_DIR, "templates")
COGNIZANT_TEMPLATE = os.path.join(TEMPLATES_DIR, "Cognizant.pptx")


# ------------------------------------------------------------
# SLIDE CLONE
# ------------------------------------------------------------
def clone_slide(prs, slide):
    new_slide = prs.slides.add_slide(slide.slide_layout)

    for shp in list(new_slide.shapes):
        new_slide.shapes._spTree.remove(shp._element)

    for shp in slide.shapes:
        new_slide.shapes._spTree.insert_element_before(
            deepcopy(shp._element), 'p:extLst'
        )

    return new_slide


# ------------------------------------------------------------
# TITLE SLIDE ONLY (UNCHANGED)
# ------------------------------------------------------------
def set_title_white_full_width(prs, slide, text):
    tf = slide.shapes.title.text_frame
    tf.clear()

    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    slide.shapes.title.left = Inches(0.5)
    slide.shapes.title.width = prs.slide_width - Inches(1)


# ------------------------------------------------------------
# CONTENT TITLE (FIXED – KEEP WHITE)
# ------------------------------------------------------------
def set_content_title(slide, title):
    if not slide.shapes.title:
        return

    tf = slide.shapes.title.text_frame
    p = tf.paragraphs[0]        # ❗ DO NOT clear
    p.text = title
    p.font.color.rgb = RGBColor(255, 255, 255)  # ✅ FORCE WHITE


# ------------------------------------------------------------
# CONTENT BODY (FIXED – BULLETS)
# ------------------------------------------------------------
def fill_content_body(slide, bullets):
    body = None
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 1:
            body = ph
            break

    if not body or not body.has_text_frame or not bullets:
        return

    tf = body.text_frame
    tf.clear()

    # ✅ FIRST BULLET — THIS IS THE KEY FIX
    tf.text = bullets[0]
    tf.paragraphs[0].level = 0
    tf.paragraphs[0].font.size = Pt(16)

    # ✅ REMAINING BULLETS
    for bullet in bullets[1:]:
        p = tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(16)

# ------------------------------------------------------------
# MAIN GENERATOR
# ------------------------------------------------------------
def generate_presentation_cognizant(payload):
    slides = payload.get("slides")
    if not slides:
        raise ValueError("No preview slides found")

    prs = Presentation(COGNIZANT_TEMPLATE)

    title_master = prs.slides[0]
    content_master = prs.slides[3]
    thankyou_master = prs.slides[-1]

    for i in reversed(range(len(prs.slides))):
        slide_id = prs.slides._sldIdLst[i].rId
        prs.part.drop_rel(slide_id)
        del prs.slides._sldIdLst[i]

    # ---------------- TITLE ----------------
    title_slide = clone_slide(prs, title_master)
    set_title_white_full_width(prs, title_slide, slides[0]["title"])

    # ---------------- CONTENT ----------------
    for slide_data in slides[1:]:
        slide = clone_slide(prs, content_master)

        set_content_title(slide, slide_data.get("title", ""))
        fill_content_body(slide, slide_data.get("bullets", []))

    # ---------------- THANK YOU ----------------
    thank_slide = clone_slide(prs, thankyou_master)
    if thank_slide.shapes.title:
        thank_slide.shapes.title.text = "Thank You"

    # ---------------- SAVE ----------------
    os.makedirs("generated", exist_ok=True)
    out = f"generated/cognizant_{uuid.uuid4().hex[:6]}.pptx"
    prs.save(out)
    return out